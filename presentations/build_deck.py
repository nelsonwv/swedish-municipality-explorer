"""
Agent 5 — builds presentations/swedish_municipality_explorer.pptx.

Run with: venv/bin/python presentations/build_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).resolve().parent / "swedish_municipality_explorer.pptx"

FONT = "Calibri"

BG = RGBColor.from_string("F5F0E8")
SURFACE = RGBColor.from_string("EDE8DC")
ACCENT = RGBColor.from_string("8B7355")
TEXT = RGBColor.from_string("2C2C2C")
TEXT_SECONDARY = RGBColor.from_string("6B6B6B")
POSITIVE = RGBColor.from_string("5B8A6F")
NEGATIVE = RGBColor.from_string("C17A5A")
BORDER = RGBColor.from_string("D4CFC6")
WHITE = RGBColor.from_string("FFFFFF")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GITHUB_URL = "https://github.com/nelsonwv/swedish-municipality-explorer"
LIVE_URL = "https://swedish-municipality-explorer.streamlit.app"
LINKEDIN_URL = "linkedin.com/in/waldeannelson/"


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def set_run(run, size=14, color=TEXT, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_paragraph(
    tf,
    text,
    size=14,
    color=TEXT,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    space_after=4,
    line_spacing=None,
    first=False,
    font=FONT,
):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic, font=font)
    return p


def add_bar(slide, left, top, width, height, color=ACCENT):
    """A plain filled rectangle with no border or shadow — used as a divider/accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_box(slide, left, top, width, height, fill=SURFACE, line_color=BORDER, line_width=1.0, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = radius
    except IndexError:
        pass
    tf = shape.text_frame
    tf.word_wrap = True
    return shape, tf


def add_info_box(
    slide,
    left,
    top,
    width,
    height,
    entries,
    fill=SURFACE,
    line_color=BORDER,
    align=PP_ALIGN.LEFT,
    line_spacing=1.05,
    vertical_anchor=MSO_ANCHOR.MIDDLE,
    margins=0.1,
    radius=0.06,
):
    """entries: list of (text, size, color, bold) tuples, one per paragraph."""
    shape, tf = add_box(slide, left, top, width, height, fill=fill, line_color=line_color, radius=radius)
    tf.margin_left = Inches(margins)
    tf.margin_right = Inches(margins)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = vertical_anchor
    for i, (text, size, color, bold) in enumerate(entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        set_run(run, size=size, color=color, bold=bold)
    return shape


def set_cell_border(cell, edges="LRTB", color="D4CFC6", width_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    tagnames = {"L": "lnL", "R": "lnR", "T": "lnT", "B": "lnB"}
    border_tags = (qn("a:lnL"), qn("a:lnR"), qn("a:lnT"), qn("a:lnB"))
    for edge in edges:
        tagname = tagnames[edge]
        tag = qn(f"a:{tagname}")
        for el in tcPr.findall(tag):
            tcPr.remove(el)
        ln = OxmlElement(f"a:{tagname}")
        ln.set("w", str(int(Pt(width_pt))))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", color)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        ln.append(dash)
        idx = sum(1 for c in tcPr if c.tag in border_tags)
        tcPr.insert(idx, ln)


def add_table(slide, left, top, width, height, headers, rows, col_widths, row_highlight=None, font_size=12, header_size=12):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    table.first_row = False
    table.horz_banding = False

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    table.rows[0].height = Inches(0.5)
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        add_paragraph(cell.text_frame, htext, size=header_size, color=WHITE, bold=True, first=True)
        set_cell_border(cell, color="8B7355")

    row_height = (height - Inches(0.5)) / len(rows)
    for i, row in enumerate(rows):
        table.rows[i + 1].height = int(row_height)
        highlight = row_highlight.get(i) if row_highlight else None
        bg = highlight if highlight else (BG if i % 2 == 0 else SURFACE)
        text_color = WHITE if highlight else TEXT
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            bold = highlight is not None and j == 0
            add_paragraph(cell.text_frame, str(val), size=font_size, color=text_color, bold=bold, first=True)
            set_cell_border(cell)
    return table


def add_title(slide, text, size=30):
    box = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12.13), Inches(0.7))
    add_paragraph(box.text_frame, text, size=size, color=ACCENT, bold=True, first=True)
    add_bar(slide, Inches(0.6), Inches(1.05), Inches(12.13), Pt(1.5), color=BORDER)
    return box


def add_footer(slide, page_num):
    box = add_textbox(slide, Inches(0.6), Inches(7.08), Inches(8), Inches(0.3))
    add_paragraph(box.text_frame, "Swedish Municipality Explorer — Waldean Nelson", size=9, color=TEXT_SECONDARY, first=True)
    box2 = add_textbox(slide, Inches(11.83), Inches(7.08), Inches(0.9), Inches(0.3))
    add_paragraph(box2.text_frame, f"{page_num} / 8", size=9, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, first=True)


# ---------------------------------------------------------------------------
# Slide 1 — Personal hook
# ---------------------------------------------------------------------------

SLIDE1_BODY = (
    "I'm a new graduate figuring out where to settle in Sweden. I found SCB data "
    "capturing employment, income, and education at municipality level — and "
    "buried in it was a region of birth dimension that nobody was talking about. "
    "As someone not born in Sweden, I couldn't ignore it. This project builds a "
    "pipeline from raw government data to an interactive dashboard that makes "
    "inequality visible at municipality level — for the first time."
)


def build_slide1(prs):
    slide = add_slide(prs)

    box = add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12.13), Inches(0.85))
    add_paragraph(box.text_frame, "Swedish Municipality Explorer", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)

    box = add_textbox(slide, Inches(0.6), Inches(1.25), Inches(12.13), Inches(0.45))
    add_paragraph(
        box.text_frame,
        "A data engineering portfolio project by Waldean Nelson",
        size=15,
        color=TEXT_SECONDARY,
        italic=True,
        align=PP_ALIGN.CENTER,
        first=True,
    )

    add_bar(slide, Inches(5.67), Inches(1.75), Inches(2.0), Pt(2), color=ACCENT)

    # Left column: personal hook narrative
    add_bar(slide, Inches(0.6), Inches(2.05), Inches(0.06), Inches(4.55), color=ACCENT)

    box = add_textbox(slide, Inches(0.9), Inches(2.05), Inches(5.6), Inches(4.55))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(tf, SLIDE1_BODY, size=16, color=TEXT, line_spacing=1.3, first=True)

    # Right column: choropleth map screenshot placeholder + caption
    box, tf = add_box(slide, Inches(6.93), Inches(2.05), Inches(5.8), Inches(3.6), fill=SURFACE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(
        tf,
        "[INSERT CHOROPLETH MAP SCREENSHOT]",
        size=14,
        color=TEXT_SECONDARY,
        italic=True,
        align=PP_ALIGN.CENTER,
        first=True,
    )

    box = add_textbox(slide, Inches(6.93), Inches(5.75), Inches(5.8), Inches(0.85))
    add_paragraph(
        box.text_frame,
        "Interactive choropleth map of all 290 Swedish municipalities colored by "
        "Opportunity Index score",
        size=12,
        color=TEXT_SECONDARY,
        italic=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
        first=True,
    )

    add_footer(slide, 1)


# ---------------------------------------------------------------------------
# Slide 2 — What it does
# ---------------------------------------------------------------------------

SLIDE2_BULLETS = [
    (
        "Opportunity scoring",
        "Scores all 290 Swedish municipalities across four dimensions — Employment, "
        "Income, Education, and Mobility — each normalized to a 0–100 percentile.",
    ),
    (
        "Custom weighting",
        "Adjustable weight sliders re-prioritize the four dimensions live, with four "
        "built-in presets: Balanced, Job seeker, Quality of life, and Family focused.",
    ),
    (
        "Integration lens",
        "A Population group filter (All residents / Born in Sweden / Born abroad) "
        "reveals integration gaps that national averages hide.",
    ),
    (
        "Interactive map",
        "An interactive choropleth map visualizes the Opportunity Index across all "
        "290 municipalities at a glance.",
    ),
]


def build_slide2(prs):
    slide = add_slide(prs)
    add_title(slide, "What it does")

    top = Inches(1.35)
    row_h = Inches(1.05)
    gap = Inches(0.18)
    for i, (lead, desc) in enumerate(SLIDE2_BULLETS):
        y = top + i * (row_h + gap)
        add_bar(slide, Inches(0.6), y + Inches(0.06), Inches(0.16), Inches(0.16), color=ACCENT)
        box = add_textbox(slide, Inches(0.95), y, Inches(6.9), row_h)
        tf = box.text_frame
        p = add_paragraph(tf, lead, size=15, color=ACCENT, bold=True, first=True, space_after=2)
        add_paragraph(tf, desc, size=13, color=TEXT, line_spacing=1.2)

    box, tf = add_box(slide, Inches(8.1), Inches(1.35), Inches(4.63), Inches(3.35), fill=SURFACE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(
        tf,
        "[INSERT SCREENSHOT: Rankings tab with Born abroad filter]",
        size=14,
        color=TEXT_SECONDARY,
        italic=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
        first=True,
    )

    box = add_textbox(slide, Inches(8.1), Inches(4.85), Inches(4.63), Inches(0.6))
    add_paragraph(
        box.text_frame,
        f"Live at: {LIVE_URL}",
        size=11,
        color=ACCENT,
        italic=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
        first=True,
    )

    add_footer(slide, 2)


# ---------------------------------------------------------------------------
# Slide 3 — Data sources
# ---------------------------------------------------------------------------

SLIDE3_HEADERS = ["Dataset", "SCB Table ID", "What it measures", "Grain"]
SLIDE3_ROWS = [
    ("Population", "BE0101A / BefolkningNy", "Total population per municipality", "290 rows — 1 per municipality"),
    ("Median income", "AA0003F / IntGr5Kom", "Median disposable income (price base amounts)", "290 rows — 1 per municipality"),
    ("Employment", "AM0207Z / RamsForvInt04N", "Employment rate, by region of birth", "870 rows — 290 munis × 3 region-of-birth groups"),
    ("Education", "UF0506B / UtbBefRegionR", "Population with post-secondary education (ISCED 5–7)", "290 rows — 1 per municipality"),
    ("Commuting", "AM0207Z / PendlingKN", "Residents commuting out of the municipality for work", "290 rows — 1 per municipality"),
]


def build_slide3(prs):
    slide = add_slide(prs)
    add_title(slide, "Data Sources (SCB PxWebApi v2)")

    add_table(
        slide,
        Inches(0.6),
        Inches(1.2),
        Inches(12.13),
        Inches(3.3),
        SLIDE3_HEADERS,
        SLIDE3_ROWS,
        col_widths=[Inches(1.6), Inches(2.4), Inches(5.83), Inches(2.3)],
        font_size=11,
    )

    box = add_textbox(slide, Inches(0.6), Inches(4.6), Inches(12.13), Inches(0.4))
    add_paragraph(
        box.text_frame,
        "All data sourced via SCB PxWebApi v2. Snapshot year: 2021. Municipality code "
        "(4-digit) is the joining key across all tables.",
        size=12,
        color=TEXT_SECONDARY,
        italic=True,
        first=True,
    )

    add_bar(slide, Inches(0.6), Inches(5.05), Inches(0.06), Inches(0.85), color=ACCENT)
    box = add_textbox(slide, Inches(0.85), Inches(5.05), Inches(11.88), Inches(0.85))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(
        tf,
        "Note: Employment is the only table with a region_of_birth dimension (total / "
        "born_in_sweden / born_abroad) — it powers the Population group filter and "
        "the integration-gap analysis on Slide 7.",
        size=12,
        color=ACCENT,
        bold=True,
        line_spacing=1.15,
        first=True,
    )

    add_bar(slide, Inches(0.6), Inches(6.0), Inches(0.06), Inches(0.95), color=BORDER)
    box = add_textbox(slide, Inches(0.85), Inches(6.0), Inches(11.88), Inches(0.95))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(
        tf,
        "Note: Employment path was replaced from LE0105Sysselsattn02N to "
        "RamsForvInt04N, and Education path was replaced from UtbSUNBefN to "
        "UtbBefRegionR after original endpoints returned 400 errors — documented "
        "in docs/agents/agent2_ingestor.md",
        size=11,
        color=TEXT_SECONDARY,
        italic=True,
        line_spacing=1.15,
        first=True,
    )

    add_footer(slide, 3)


# ---------------------------------------------------------------------------
# Slide 4 — The data pipeline
# ---------------------------------------------------------------------------

SLIDE4_BOXES = [
    ("SCB PxWebApi v2", "5 datasets · 2021 snapshot · 290 municipalities"),
    ("Python Ingestor", "SCBClient + SnowflakeLoader — rate-limited, cell-aware chunking"),
    ("Snowflake RAW", "5 raw tables · ~2,030 rows loaded"),
    ("dbt Staging", "5 stg_* models — cleaned & standardized columns"),
    ("dbt Intermediate + Marts", "int_municipality_profiles → dim_municipality + fct_opportunity_scores"),
    ("Streamlit Dashboard", "Rankings · Compare · Map · About"),
]


def build_slide4(prs):
    slide = add_slide(prs)
    add_title(slide, "The Data Pipeline")

    n = len(SLIDE4_BOXES)
    arrow_w = Inches(0.32)
    total_arrows = arrow_w * (n - 1)
    box_w = (Inches(12.13) - total_arrows) // n
    top = Inches(2.1)
    box_h = Inches(2.3)

    x = Inches(0.6)
    for i, (name, annotation) in enumerate(SLIDE4_BOXES):
        add_info_box(
            slide,
            x,
            top,
            box_w,
            box_h,
            entries=[
                (name, 13, ACCENT, True),
                (annotation, 10, TEXT, False),
            ],
            fill=SURFACE,
            align=PP_ALIGN.CENTER,
            line_spacing=1.1,
        )
        x += box_w
        if i < n - 1:
            box = add_textbox(slide, x, top, arrow_w, box_h)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            add_paragraph(tf, "→", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)
            x += arrow_w

    add_bar(slide, Inches(0.6), Inches(5.05), Inches(12.13), Pt(1.5), color=BORDER)

    box = add_textbox(slide, Inches(0.6), Inches(5.25), Inches(12.13), Inches(0.6))
    add_paragraph(
        box.text_frame,
        "53 dbt tests  •  Full documentation  •  Version controlled on GitHub  "
        "•  Agent-documented architecture",
        size=15,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
        first=True,
    )

    add_footer(slide, 4)


# ---------------------------------------------------------------------------
# Slide 5 — Data model: lineage + star schema
# ---------------------------------------------------------------------------

LINEAGE_LAYERS = [
    ("RAW", ["raw_population", "raw_income", "raw_employment ← (key: has region_of_birth)", "raw_education", "raw_commuting"]),
    ("STAGING", ["stg_population", "stg_income", "stg_employment", "stg_education", "stg_commuting"]),
    ("INTERMEDIATE", ["int_municipality_profiles (290 rows)", "Calculated: employment_gap, pct_post_secondary, out_commute_rate"]),
    ("MARTS", ["dim_municipality", "fct_opportunity_scores"]),
]

FACT_COLUMNS = [
    "fct_opportunity_scores",
    "• municipality_code",
    "• region_of_birth",
    "• income_score",
    "• employment_score",
    "• education_score",
    "• mobility_score",
    "• year",
]

DIM_COLUMNS = [
    "dim_municipality",
    "• municipality_code (PK)",
    "• municipality_name",
    "• county_name",
    "• population",
    "• population_category",
    "• year",
]


def build_slide5(prs):
    slide = add_slide(prs)
    add_title(slide, "Data Model — From Raw to Star Schema")

    add_bar(slide, Inches(6.8), Inches(1.3), Pt(1.5), Inches(5.5), color=BORDER)

    box = add_textbox(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.4))
    add_paragraph(box.text_frame, "Full Lineage (dbt models)", size=15, color=ACCENT, bold=True, first=True)

    box = add_textbox(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.4))
    add_paragraph(box.text_frame, "Star Schema (MARTS)", size=15, color=ACCENT, bold=True, first=True)

    n = len(LINEAGE_LAYERS)
    arrow_h = Inches(0.18)
    layer_h = Inches(1.05)
    top = Inches(1.8)
    for i, (layer_name, tables) in enumerate(LINEAGE_LAYERS):
        entries = [(layer_name, 11, ACCENT, True)] + [(t, 9, TEXT, False) for t in tables]
        add_info_box(
            slide,
            Inches(0.6),
            top,
            Inches(6.0),
            layer_h,
            entries=entries,
            fill=SURFACE,
            line_spacing=1.05,
            margins=0.12,
        )
        top += layer_h
        if i < n - 1:
            box = add_textbox(slide, Inches(0.6), top, Inches(6.0), arrow_h)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            add_paragraph(tf, "↓", size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
            top += arrow_h

    fact_entries = [(FACT_COLUMNS[0], 12, ACCENT, True)] + [(c, 10, TEXT, False) for c in FACT_COLUMNS[1:]]
    fact_box = add_info_box(
        slide,
        Inches(7.1),
        Inches(2.1),
        Inches(2.9),
        Inches(2.0),
        entries=fact_entries,
        fill=SURFACE,
        line_spacing=1.08,
        align=PP_ALIGN.LEFT,
    )

    dim_entries = [(DIM_COLUMNS[0], 12, ACCENT, True)] + [(c, 10, TEXT, False) for c in DIM_COLUMNS[1:]]
    dim_box = add_info_box(
        slide,
        Inches(10.4),
        Inches(2.25),
        Inches(2.4),
        Inches(1.7),
        entries=dim_entries,
        fill=SURFACE,
        line_spacing=1.08,
        align=PP_ALIGN.LEFT,
    )

    connector_y = Inches(3.1)
    line = add_bar(slide, Inches(10.0), connector_y, Inches(0.4), Pt(1.5), color=ACCENT)

    box = add_textbox(slide, Inches(7.1), Inches(4.4), Inches(5.7), Inches(0.9))
    add_paragraph(
        box.text_frame,
        "Each municipality has 3 fct_opportunity_scores rows — one per "
        "region_of_birth value (total / born_in_sweden / born_abroad) — joined to a "
        "single dim_municipality row on (municipality_code, year).",
        size=11,
        color=TEXT_SECONDARY,
        italic=True,
        line_spacing=1.2,
        first=True,
    )

    add_footer(slide, 5)


# ---------------------------------------------------------------------------
# Slide 6 — How the Opportunity Index works
# ---------------------------------------------------------------------------

SLIDE6_CARDS = [
    (
        "Employment",
        "Min-max normalized within each region-of-birth group — reflects a "
        "municipality's standing among peers for that specific population group "
        "(total / born in Sweden / born abroad).",
    ),
    (
        "Income",
        "Min-max normalized median income across all 290 municipalities — higher "
        "income means a higher score.",
    ),
    (
        "Education",
        "Share of the population (16–95+) with post-secondary education "
        "(ISCED 5–7), min-max normalized across all 290 municipalities.",
    ),
    (
        "Mobility",
        "Inverse min-max of the out-commuting rate — municipalities where fewer "
        "residents commute elsewhere for work score higher.",
    ),
]

SLIDE6_FORMULA = (
    "opportunity_index = (\n"
    "    employment_score * w_employment\n"
    "  + income_score      * w_income\n"
    "  + education_score   * w_education\n"
    "  + mobility_score    * w_mobility\n"
    ") / sum(all weights)"
)

SLIDE6_CALLOUT = (
    "Think of each score as a percentile — a score of 75 means this municipality "
    "outperforms 75% of Sweden's 290 municipalities on that dimension. Danderyd "
    "scores 100 on income (highest median in Sweden) but only 56 on employment — "
    "its large wealthy retired population aren't in the workforce. The index "
    "captures that trade-off in one number."
)


def build_slide6(prs):
    slide = add_slide(prs)
    add_title(slide, "How the Opportunity Index Works")

    card_w = Inches(5.9)
    card_h = Inches(1.3)
    gap = Inches(0.33)
    positions = [
        (Inches(0.6), Inches(1.25)),
        (Inches(0.6) + card_w + gap, Inches(1.25)),
        (Inches(0.6), Inches(1.25) + card_h + Inches(0.1)),
        (Inches(0.6) + card_w + gap, Inches(1.25) + card_h + Inches(0.1)),
    ]
    for (left, top), (name, desc) in zip(positions, SLIDE6_CARDS):
        add_info_box(
            slide,
            left,
            top,
            card_w,
            card_h,
            entries=[
                (name, 14, ACCENT, True),
                (desc, 11, TEXT, False),
            ],
            fill=SURFACE,
            line_spacing=1.15,
            vertical_anchor=MSO_ANCHOR.TOP,
            margins=0.14,
        )

    code_box, tf = add_box(slide, Inches(0.6), Inches(4.15), Inches(12.13), Inches(1.5), fill=SURFACE)
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(SLIDE6_FORMULA.split("\n")):
        add_paragraph(tf, line, size=13, color=TEXT, font="Consolas", first=(i == 0), space_after=0, line_spacing=1.15)

    box, tf = add_box(slide, Inches(0.6), Inches(5.85), Inches(12.13), Inches(1.1), fill=ACCENT, line_color=ACCENT)
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(tf, SLIDE6_CALLOUT, size=13, color=WHITE, line_spacing=1.25, first=True)

    add_footer(slide, 6)


# ---------------------------------------------------------------------------
# Slide 7 — The gap that matters
# ---------------------------------------------------------------------------

SLIDE7_OPENING = (
    "Switch the Population group filter from 'All residents' to 'Born abroad' "
    "and watch the rankings reshuffle — this is the moment the dashboard was "
    "built for."
)

SLIDE7_HEADERS = ["Municipality", "Born in Sweden rank", "Born abroad rank", "Change"]
SLIDE7_ROWS = [
    ("Sollentuna", "#15", "#10", "▲ +5"),
    ("Österåker", "#16", "#12", "▲ +4"),
    ("Ekerö", "#8", "#7", "▲ +1"),
    ("Göteborg", "#9", "#8", "▲ +1"),
    ("Härryda", "#10", "#11", "→ stable"),
    ("Solna", "#7", "#9", "▼ -2"),
    ("Stockholm", "#2", "#3", "▼ -1"),
]
SLIDE7_HIGHLIGHT = {0: POSITIVE, 1: POSITIVE, 5: NEGATIVE, 6: NEGATIVE}

SLIDE7_BOTTOM = (
    "For the large majority of Sweden's 290 municipalities, overall rank and "
    "born-abroad rank are nearly identical — meaning opportunity is largely "
    "independent of birthplace. The exceptions are where this dashboard adds "
    "value: invisible in national averages, visible only at municipality level."
)


def build_slide7(prs):
    slide = add_slide(prs)
    add_title(slide, "The Gap That Matters")

    box = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12.13), Inches(0.85))
    add_paragraph(box.text_frame, SLIDE7_OPENING, size=19, color=ACCENT, bold=True, line_spacing=1.2, first=True)

    add_table(
        slide,
        Inches(0.6),
        Inches(2.15),
        Inches(12.13),
        Inches(2.85),
        SLIDE7_HEADERS,
        SLIDE7_ROWS,
        col_widths=[Inches(3.0), Inches(2.7), Inches(2.7), Inches(3.73)],
        row_highlight=SLIDE7_HIGHLIGHT,
        font_size=13,
    )

    box_w = Inches(5.9)
    add_info_box(
        slide,
        Inches(0.6),
        Inches(5.2),
        box_w,
        Inches(1.0),
        entries=[
            ("KEY INSIGHT", 10, POSITIVE, True),
            (
                "Sollentuna jumps from #15 to #10 (+5) for residents born abroad — "
                "its employment outcomes for foreign-born residents outperform its "
                "overall ranking, suggesting strong labor-market integration.",
                11,
                TEXT,
                False,
            ),
        ],
        fill=SURFACE,
        line_spacing=1.1,
        vertical_anchor=MSO_ANCHOR.TOP,
        margins=0.14,
    )
    add_info_box(
        slide,
        Inches(0.6) + box_w + Inches(0.33),
        Inches(5.2),
        box_w,
        Inches(1.0),
        entries=[
            ("KEY INSIGHT", 10, NEGATIVE, True),
            (
                "Stockholm slips from #2 to #3 (-1) for residents born abroad. Even "
                "Sweden's capital — and the country's #2 municipality overall — "
                "shows a measurable employment gap by region of birth.",
                11,
                TEXT,
                False,
            ),
        ],
        fill=SURFACE,
        line_spacing=1.1,
        vertical_anchor=MSO_ANCHOR.TOP,
        margins=0.14,
    )

    box, tf = add_box(slide, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.7), fill=SURFACE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_paragraph(tf, SLIDE7_BOTTOM, size=11, color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.2, first=True)

    add_footer(slide, 7)


# ---------------------------------------------------------------------------
# Slide 8 — Challenges, future improvements & closing
# ---------------------------------------------------------------------------

SLIDE8_CHALLENGES = [
    "SCB API endpoint changes — 2 of the 5 planned table paths returned "
    "400 Bad Request and had to be replaced with current equivalents "
    "(employment, education) before ingestion could run.",
    "A 2019 methodology break in SCB's employment statistics meant older series "
    "weren't comparable to the 2021 snapshot, ruling out a simple historical "
    "extension.",
    "Commuting data was originally planned from Trafikanalys, but that source is "
    "county-level only — replaced with SCB's municipality-level commuting table "
    "(PendlingKN) instead.",
    "SCB's PxWebApi enforces a 150,000-cell limit per query — the ingestor "
    "includes cell-count-aware chunking logic for any future larger queries.",
]

SLIDE8_IMPROVEMENTS = [
    "Automate the pipeline with GitHub Actions for scheduled refreshes as new SCB "
    "data is published.",
    "Add a 5th dimension — housing cost/affordability — sourced from Booli or "
    "Hemnet listing data.",
    "Extend to historical trending (2015–2023) once the 2019 methodology break is "
    "reconciled.",
    "Incorporate school performance data (e.g. Skolverket) as an additional "
    "opportunity dimension.",
    "Migrate from Snowflake to DuckDB for a lighter-weight, fully local/open-source "
    "deployment.",
    "Add statistical significance testing to the rank-change analysis on Slide 7.",
]


def build_slide8(prs):
    slide = add_slide(prs)
    add_title(slide, "Challenges, Future Improvements & Closing")

    add_bar(slide, Inches(6.8), Inches(1.3), Pt(1.5), Inches(3.6), color=BORDER)

    box = add_textbox(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.4))
    add_paragraph(box.text_frame, "Challenges encountered", size=15, color=ACCENT, bold=True, first=True)

    box = add_textbox(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.4))
    add_paragraph(box.text_frame, "Future improvements", size=15, color=ACCENT, bold=True, first=True)

    box = add_textbox(slide, Inches(0.6), Inches(1.8), Inches(6.0), Inches(3.3))
    tf = box.text_frame
    for i, item in enumerate(SLIDE8_CHALLENGES):
        add_paragraph(tf, f"•  {item}", size=12, color=TEXT, line_spacing=1.15, space_after=10, first=(i == 0))

    box = add_textbox(slide, Inches(7.0), Inches(1.8), Inches(6.0), Inches(3.3))
    tf = box.text_frame
    for i, item in enumerate(SLIDE8_IMPROVEMENTS):
        add_paragraph(tf, f"•  {item}", size=12, color=TEXT, line_spacing=1.15, space_after=8, first=(i == 0))

    closing_top = Inches(5.25)
    closing_h = Inches(1.55)
    card_w = Inches(3.91)
    gap = Inches(0.2)

    add_info_box(
        slide,
        Inches(0.6),
        closing_top,
        card_w,
        closing_h,
        entries=[
            ("Live Dashboard", 13, ACCENT, True),
            (LIVE_URL, 12, TEXT, False),
        ],
        fill=SURFACE,
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )
    add_info_box(
        slide,
        Inches(0.6) + card_w + gap,
        closing_top,
        card_w,
        closing_h,
        entries=[
            ("GitHub", 13, ACCENT, True),
            (GITHUB_URL, 12, TEXT, False),
        ],
        fill=SURFACE,
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )
    add_info_box(
        slide,
        Inches(0.6) + 2 * (card_w + gap),
        closing_top,
        card_w,
        closing_h,
        entries=[
            ("Connect", 13, ACCENT, True),
            (LINKEDIN_URL, 12, TEXT, False),
        ],
        fill=SURFACE,
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )

    add_footer(slide, 8)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main():
    prs = new_presentation()
    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    build_slide8(prs)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
